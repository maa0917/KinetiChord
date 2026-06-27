import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # カラー定義
    BG_COLOR = RGBColor(10, 4, 20)      # 極暗紫
    CARD_BG = RGBColor(22, 11, 38)       # プレミアムグラスカード風紫
    CYAN = RGBColor(0, 240, 255)         # ネオンシアン
    PINK = RGBColor(255, 42, 133)        # ネオンピンク
    GOLD = RGBColor(255, 215, 0)         # ゴールド
    WHITE = RGBColor(255, 255, 255)      # 純白
    MUTED = RGBColor(159, 143, 179)      # ミュートグレー紫
    BLUE = RGBColor(52, 152, 219)        # サッドブルー
    RED = RGBColor(255, 51, 51)          # アングリーレッド
    GREEN = RGBColor(0, 255, 170)        # ニュートラルグリーン

    blank_layout = prs.slide_layouts[6]

    def apply_background(slide):
        # 背景全画面長方形
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

    def add_header(slide, title_text, category_text, page_str):
        # ヘッダー境界線
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.2), Inches(11.833), Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(40, 25, 65)
        line.line.fill.background()
        
        # 左上ロゴ
        logo_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(5.0), Inches(0.8))
        tf = logo_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "✦ KinetiChord"
        p.font.name = 'Outfit'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PINK
        
        # 右上カテゴリタグ
        tag_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.4), Inches(4.083), Inches(0.8))
        tf_tag = tag_box.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = category_text.upper()
        p_tag.alignment = PP_ALIGN.RIGHT
        p_tag.font.name = 'Outfit'
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = CYAN
        
        # スライド大タイトル
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.833), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Noto Sans JP'
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        # フッター境界線
        f_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.7), Inches(11.833), Inches(0.015))
        f_line.fill.solid()
        f_line.fill.fore_color.rgb = RGBColor(30, 15, 50)
        f_line.line.fill.background()
        
        # フッタークレジット
        f_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.75), Inches(6.0), Inches(0.5))
        p_f = f_box.text_frame.paragraphs[0]
        p_f.text = "Gemini AI Hackathon 2026 // Presentation Deck"
        p_f.font.name = 'Outfit'
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = MUTED
        
        # フッターページ番号
        p_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.75), Inches(2.083), Inches(0.5))
        p_p = p_box.text_frame.paragraphs[0]
        p_p.text = page_str
        p_p.alignment = PP_ALIGN.RIGHT
        p_p.font.name = 'Outfit'
        p_p.font.size = Pt(11)
        p_p.font.bold = True
        p_p.font.color.rgb = WHITE

    def create_card(slide, left, top, width, height, title, items, color_title=WHITE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(60, 30, 90)
        card.line.width = Pt(1.5)
        
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = f"✦ {title}"
        p_title.font.name = 'Noto Sans JP'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = color_title
        p_title.space_after = Pt(12)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.name = 'Noto Sans JP'
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

    # ==========================================================
    # SLIDE 1: COVER
    # ==========================================================
    s1 = prs.slides.add_slide(blank_layout)
    apply_background(s1)
    
    # 巨大メインタイトル
    m_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5))
    tf = m_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "KinetiChord"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(84)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "身体と感情のダイナミズムが織りなす、極限の音響＆ビジュアル共鳴アート"
    p2.font.name = 'Noto Sans JP'
    p2.font.size = Pt(22)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)

    # デコレーションライン
    dec = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.666), Inches(4.4), Inches(4.0), Inches(0.04))
    dec.fill.solid()
    dec.fill.fore_color.rgb = PINK
    dec.line.fill.background()

    # メタ情報
    meta_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(2.0))
    tf_m = meta_box.text_frame
    p_m1 = tf_m.paragraphs[0]
    p_m1.text = "Developer: maa0917 // Gemini AI Hackathon 2026"
    p_m1.font.name = 'Outfit'
    p_m1.font.size = Pt(14)
    p_m1.font.color.rgb = WHITE
    p_m1.alignment = PP_ALIGN.CENTER
    
    p_m2 = tf_m.add_paragraph()
    p_m2.text = "synth-flow-1038264562195.us-central1.run.app"
    p_m2.font.name = 'Outfit'
    p_m2.font.size = Pt(12)
    p_m2.font.color.rgb = CYAN
    p_m2.alignment = PP_ALIGN.CENTER
    p_m2.space_before = Pt(8)

    # ==========================================================
    # SLIDE 2: BACKGROUND & CONCEPT
    # ==========================================================
    s2 = prs.slides.add_slide(blank_layout)
    apply_background(s2)
    add_header(s2, "開発背景とコア・コンセプト", "01. Background & Concept", "Page 02")
    
    create_card(s2, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "現代のデジタルアートの課題", [
        "「静観」から抜け出せない体験: 画面の向こう側の鑑賞にとどまり、身体感覚と切り離されている。",
        "双方向性の希薄さ: 単一のマウスクリックや固定ボタン操作のみといった、単調なインタラクションの限界。",
        "感情の未反映: ユーザーのその瞬間の感情や心のゆらぎが、表現の挙動にダイレクトに反映されない。"
    ], PINK)
    
    create_card(s2, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "KinetiChordの解答", [
        "身体を「楽器」に、表情を「光」に変え、空間全体を共鳴させる新機軸のデジタルアート。",
        "マルチモーダル融合: 表情分析、手の重心追跡、全体のモーション強度、マイク音圧解析の完全ブレンド。",
        "音と光をまとう超感覚: 体の移動速度や激しさに応じて、リアルタイムに粒子とシンセ音響が爆発・変調。"
    ], CYAN)

    # ==========================================================
    # SLIDE 3: SYSTEM OVERVIEW
    # ==========================================================
    s3 = prs.slides.add_slide(blank_layout)
    apply_background(s3)
    add_header(s3, "入力から出力へのシームレスな共鳴ロジック", "02. System Architecture", "Page 03")
    
    # 3カラム
    create_card(s3, Inches(0.75), Inches(2.3), Inches(3.6), Inches(4.0), "1. センシング (入力)", [
        "表情/感情分析: face-api.jsを用いて100ms周期で顔の喜怒哀楽スコアを検出。",
        "身体モーション: 手の重心座標や全体のフレーム差分ピクセル量を計算。",
        "マイク音圧: WebAudio APIを通じて音声ビートを完全追跡。"
    ], PINK)
    
    create_card(s3, Inches(4.866), Inches(2.3), Inches(3.6), Inches(4.0), "2. 共鳴処理 (コアコア)", [
        "3Dパース射影: 3D座標データをリアルタイムに2D画面へとカメラ射影演算。",
        "Lerpスムージング: 不自然なカクつきを根絶し、追従遅延ゼロの滑らかな補間処理を実行。",
        "音響処理グラフ: 感情ごとの特殊エフェクトノードを動的に結線。"
    ], CYAN)
    
    create_card(s3, Inches(8.983), Inches(2.3), Inches(3.6), Inches(4.0), "3. アウトプット (出力)", [
        "3D Canvas: 420個の密な粒子、3Dテーパーシリンダーによる立体肉厚スケルトン。",
        "立体音響シンセ: 手の位置に完全に同期するステレオ定位パンニング処理。",
        "極限音響: 脳を揺さぶるエフェクト。"
    ], GOLD)

    # ==========================================================
    # SLIDE 4: CORE FEATURE 1
    # ==========================================================
    s4 = prs.slides.add_slide(blank_layout)
    apply_background(s4)
    add_header(s4, "身体を完璧にトレースする「キネティック・スケルトン」", "03. Core Feature I", "Page 04")
    
    create_card(s4, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "二律駆動 (Dual-Mode) アバター", [
        "リアルタイム追従: 顔座標と手の重心をミリ秒単位でスムージング追跡。",
        "自律遊泳ダンス回帰: カメラ検出から外れた瞬間に、呼吸するような2重調和振動（∞字ウェーブ）による自律遊泳ポーズへ移行。",
        "フリーズ根絶: 人がいない時も生き生きとした美しい浮遊感を表現。"
    ], CYAN)
    
    create_card(s4, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "3Dシリンダー投影 ＆ 幾何数学", [
        "脚部3Dボリューム復元: 股関節から足首までの骨格を3Dテーパーシリンダー（肉厚な幾何シリンダー）として数理構築。",
        "シネマティック旋回アングル: 自動でゆっくりとアバターが左右に自転、上下にたゆたう3Dパースペクティブ投影カメラシステム。",
        "完全な奥行き表現: ドラッグ＆タッチ操作で全方位360度から鑑賞・回転可能。"
    ], PINK)

    # ==========================================================
    # SLIDE 5: CORE FEATURE 2
    # ==========================================================
    s5 = prs.slides.add_slide(blank_layout)
    apply_background(s5)
    add_header(s5, "動きに完全に連動するダイナミック・パーティクル", "04. Core Feature II", "Page 05")
    
    create_card(s5, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "手速度追跡 ＆ 物理軌跡バースト", [
        "速度ベクトルのリアルタイム算出: 手（smoothLeftHand/RightHand）のフレーム間移動差分から瞬時速度を検出。",
        "物理軌跡バースト: 手の移動速度が一定値を超えた瞬間、進行方向の慣性を引き継いだ3〜8個の軌跡バースト粒子を鋭く射出。",
        "手から光を放つような爽快な直感フィードバック体験を実現。"
    ], CYAN)
    
    create_card(s5, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "全身モーションスパーク ＆ 星座脈動", [
        "全身モーションスパーク: ウェブカメラ全体の動きの激しさに応じて、アバター骨格から瞬く「虹色の極小スパーク」を毎フレーム大放出。",
        "ダイナミック星座ネットワーク: 人体トポロジー（areAdjacentBodyParts）に基づき不要ノイズ線を完全遮断。激しい動作時はラインが太く脈動し、最大接続距離が自動拡張。"
    ], PINK)

    # ==========================================================
    # SLIDE 6: CORE FEATURE 3
    # ==========================================================
    s6 = prs.slides.add_slide(blank_layout)
    apply_background(s6)
    add_header(s6, "感情を100%増幅・反映する極限音響システム", "05. Core Feature III", "Page 06")
    
    # 2カラムで、感情ごとの音響モジュレーション
    create_card(s6, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "😊 喜び (Happy) ＆ 😢 哀しみ (Sad)", [
        "喜び: フィルターカットオフを9,500Hzまで完全解放。120ms周期でコード音がオクターブ上（+12/+24半音）へスイープする流麗なアルペジエーター。",
        "哀しみ: 深海フェードアウト音量85%カット。ローパスを150Hzの極限深海マッフル状態にし、88%帰還率のディレイによる底知れぬ哀愁のエコー。"
    ], GOLD)
    
    create_card(s6, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "😡 怒り (Angry) ＆ 😮 驚き (Surprised)", [
        "怒り: クリーン経路を完全遮断。WaveShaperによるディストーションの歪み量を極限（350）にブースト。マイク音圧連動の爆音ゲインと高周波悲鳴ピッチジッター。",
        "驚き: 15Hz超ハイスピードステレオ往復パンニングが脳髄を直撃。0.05秒の極小スラップバックディレイによる金属的なショートエコー。"
    ], RED)

    # ==========================================================
    # SLIDE 7: TECH STACK
    # ==========================================================
    s7 = prs.slides.add_slide(blank_layout)
    apply_background(s7)
    add_header(s7, "徹底したPure Vanilla実装と極限パフォーマンス最適化", "06. Tech Stack", "Page 07")
    
    create_card(s7, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "技術スタック一覧", [
        "Pure Vanilla Web: 重厚なReact等のUIフレームワークを排除。HTML/CSS/Vanilla JSによる瞬時ロードと超レスポンシブな描画を実現。",
        "face-api.js (TensorFlow.js): 軽量な表情分析モデル。約5MBのデータは起動と同時に非同期読み込みされ、UI操作を一切阻害しません。",
        "Web Audio API: シンセ合成およびエフェクト結線をWebブラウザ上で完結。"
    ], CYAN)
    
    create_card(s7, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "60FPSを描画し続けるチューニング技術", [
        "超軽量差分バッファ: モーション判定をわずか「48x36」の極小解像度で行い、CPU/GPU消費を事実上ゼロへクランプしながら手の重心を追従。",
        "厳格なメモリ・アロケーション: パーティクル最大容量を「650」にハードクランプ。寿命を迎えた粒子を即時GCへ。モバイルでもヌルヌル動作。"
    ], PINK)

    # ==========================================================
    # SLIDE 8: DEPLOYMENT & VERIFICATION
    # ==========================================================
    s8 = prs.slides.add_slide(blank_layout)
    apply_background(s8)
    add_header(s8, "Docker ＆ Google Cloud Run によるサーバーレスデプロイ", "07. Deployment", "Page 08")
    
    create_card(s8, Inches(0.75), Inches(2.3), Inches(5.6), Inches(4.0), "デプロイ構成の真髄", [
        "Dockerコンテナ化: Nginxによる極小のAlpine Linuxイメージを定義し、軽量ポータビリティを担保。",
        "Google Cloud Run: サーバーレス環境。常時起動（ミニマムインスタンス）設定により、起動遅延（コールドスタート）を完璧に根絶。",
        "完全常時HTTPS: カメラ・マイクのブラウザセキュリティポリシーに完全対応するためのセキュアなSSL通信の担保。"
    ], CYAN)
    
    create_card(s8, Inches(6.983), Inches(2.3), Inches(5.6), Inches(4.0), "動作検証テストガイド", [
        "Step 1: 「Start」をクリックし、マイクとカメラフィードを承認。",
        "Step 2 (マイク拍動): 手拍子や声で白銀粒子が一斉バーストするのを確認。",
        "Step 3 (感情大爆発): 顔の表情を「喜び」や「怒り」に変え、関節からゴールドや真紅の粒子が30〜50個規模で勢いよく大爆発するのを確認。",
        "Step 4 (手スワイプ): 手を高速スワイプし、進行方向に物理弾が射出。"
    ], PINK)

    # ==========================================================
    # SLIDE 9: FUTURE VISION
    # ==========================================================
    s9 = prs.slides.add_slide(blank_layout)
    apply_background(s9)
    add_header(s9, "KinetiChordが切り拓く将来の展望と可能性", "08. Future Vision", "Page 09")
    
    create_card(s9, Inches(0.75), Inches(3.6), Inches(3.6), Inches(2.8), "1. XR空間展開 (VR/AR)", [
        "Apple Vision Pro等の3Dハンドトラッキングセンサーと統合。",
        "全方位の3D空間全体へ手から数万個の音響粒子を放つ、メタバース・アートへ進化。"
    ], CYAN)
    
    create_card(s9, Inches(4.866), Inches(3.6), Inches(3.6), Inches(2.8), "2. 外部DAW連携 (MIDI)", [
        "Web MIDI API経由でAbleton LiveやLogic Pro等のプロ音楽ソフトと直結。",
        "ダンサーの身体そのものがシンセを狂暴に変調させる、身体シンセ。"
    ], PINK)
    
    create_card(s9, Inches(8.983), Inches(3.6), Inches(3.6), Inches(2.8), "3. ライブプロジェクション", [
        "高輝度プロジェクターを使用し、ダンスステージに3Dアバターをリアルタイム投影。",
        "演者の身体と劇的なリアルタイム視覚エフェクトが寸分の遅れもなく共鳴。"
    ], GOLD)

    # ==========================================================
    # SLIDE 10: CONCLUSION
    # ==========================================================
    s10 = prs.slides.add_slide(blank_layout)
    apply_background(s10)
    
    m_box10 = s10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5))
    tf10 = m_box10.text_frame
    tf10.word_wrap = True
    p10_1 = tf10.paragraphs[0]
    p10_1.text = "身体は楽器になり、感情は光となる。"
    p10_1.font.name = 'Noto Sans JP'
    p10_1.font.size = Pt(36)
    p10_1.font.bold = True
    p10_1.font.color.rgb = CYAN
    p10_1.alignment = PP_ALIGN.CENTER
    p10_1.space_after = Pt(20)
    
    p10_2 = tf10.add_paragraph()
    p10_2.text = "Thank You."
    p10_2.font.name = 'Outfit'
    p10_2.font.size = Pt(72)
    p10_2.font.bold = True
    p10_2.font.color.rgb = WHITE
    p10_2.alignment = PP_ALIGN.CENTER
    
    p10_3 = tf10.add_paragraph()
    p10_3.text = "ご清聴ありがとうございました。"
    p10_3.font.name = 'Noto Sans JP'
    p10_3.font.size = Pt(20)
    p10_3.font.color.rgb = MUTED
    p10_3.alignment = PP_ALIGN.CENTER
    p10_3.space_before = Pt(10)

    # デコレーションライン
    dec10 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.666), Inches(5.0), Inches(4.0), Inches(0.04))
    dec10.fill.solid()
    dec10.fill.fore_color.rgb = PINK
    dec10.line.fill.background()

    meta10 = s10.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.333), Inches(1.5))
    tf_m10 = meta10.text_frame
    p_m10_1 = tf_m10.paragraphs[0]
    p_m10_1.text = "Project Repository: github.com/maa0917/KinetiChord"
    p_m10_1.font.name = 'Outfit'
    p_m10_1.font.size = Pt(14)
    p_m10_1.font.color.rgb = WHITE
    p_m10_1.alignment = PP_ALIGN.CENTER

    prs.save("/Users/matsuhashitatsuya/WorkSpace/Gemini AI Hackathon/KinetiChord_Slides.pptx")
    print("Success: KinetiChord_Slides.pptx created successfully!")

if __name__ == "__main__":
    create_deck()
